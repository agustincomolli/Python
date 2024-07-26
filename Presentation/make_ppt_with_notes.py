from docx import Document
from pptx import Presentation
from pptx.util import Pt

# Leer el documento de Word
doc_path = "./Charla sobre Lo Sagrado.docx"
doc = Document(doc_path)

# Función para extraer el contenido y la página de cada párrafo
def extract_content_and_pages(doc):
    content_pages = []
    current_page = 1
    for para in doc.paragraphs:
        if para.style.name == 'Heading 1' or para.style.name == 'Heading 2':
            current_page += 1
        content_pages.append((para.text, current_page))
    return content_pages

# Extraer el contenido y las páginas del documento
content_pages = extract_content_and_pages(doc)

# Crear una nueva presentación
prs = Presentation()

# Definir las diapositivas y sus contenidos junto con las páginas de origen
slides_content_with_pages = [
    ("El Sentido de lo Sagrado", "La Fe Católica frente a la Indiferencia Moderna", 1),
    ("Lo Tradicional y la Biblia", "La Iglesia vive de la Biblia y de la Tradición.\nImportancia igual de ambos según el Vaticano II (DV 9).", 1),
    ("La Iglesia y la Tradición", "La Biblia nace de la Tradición.\nSin la luz de la Tradición, la Biblia no nos valdría para nada.", 1),
    ("Lo Sagrado", "Definición de lo Sagrado\nAlgo apartado o consagrado para un propósito espiritual.\nConsiderado digno de respeto o veneración especial.", 2),
    ("Lo Sagrado vs. Lo Profano", "Distinción entre lo Sagrado y lo Profano\nLo sagrado está conectado con la divinidad.\nLo profano refiere a lo ordinario y cotidiano.", 2),
    ("Manifestaciones de lo Sagrado", "Diversas Formas de lo Sagrado\n- Objetos Sagrados: libros religiosos, reliquias, etc.\n- Lugares Sagrados: templos, santuarios, etc.\n- Personas Sagradas: sacerdotes, monjes, etc.\n- Textos Sagrados: Biblia, Corán, etc.\n- Tiempo Sagrado: días santos, festivales, etc.\n- Acciones Sagradas: rituales, ceremonias, etc.", 2),
    ("Continuidad de lo Sagrado", "De lo Pagano a lo Cristiano\nLa gracia perfecciona la naturaleza.\nContinuidad desde religiosidades naturales hasta la adoración cristiana.", 3),
    ("Lo Sagrado Judío", "La Sacralidad en el Judaísmo\nOrden de sacralidades: fiestas, sacerdocio, lugares, etc.\nIsrael como pueblo sagrado.", 3),
    ("Lo Sagrado Cristiano", "Cristo y la Sacralidad Cristiana\nCristo como sagrado absoluto.\nLa Iglesia como 'sacramento universal de salvación'.", 3),
    ("Teología de lo Sagrado", "Definición Teológica\nJesucristo es sagrado por su humanidad.\nDistinción entre santo y sagrado.", 4),
    ("Secularización y Desacralización", "Impacto Moderno en lo Sagrado\nPérdida de sensibilidad para lo sagrado en Occidente.\nConsecuencias espirituales de esta pérdida.", 4),
    ("Espiritualidad Cristiana", "Amor a lo Sagrado\nImportancia del orden sacral en la espiritualidad católica.\nBúsqueda del Santo en las cosas sagradas de la Iglesia.", 4),
    ("Tendencia de lo Sagrado a la Manifestación", "Visibilidad de lo Sagrado\nLo sagrado como signo claro y visible.\nEjemplos: templos, vestimentas religiosas, campanas, canto religioso.", 5),
    ("Conclusión", "La Relevancia de lo Sagrado Hoy\nImportancia de mantener y respetar las formas sagradas.\nPapel de lo sagrado en la identidad y misión de la Iglesia.", 5)
]

# Crear y añadir las diapositivas con notas
for title, content, page in slides_content_with_pages:
    slide_layout = prs.slide_layouts[1]  # Usar el layout de título y contenido
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    content_placeholder.text = content
    
    # Añadir notas a la diapositiva
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    p = text_frame.add_paragraph()
    p.text = f"Contenido extraído de la página {page} del documento de Word."
    p.font.size = Pt(12)

# Guardar la presentación en un archivo
pptx_file_with_notes = "./Charla_sobre_Lo_Sagrado_con_notas.pptx"
prs.save(pptx_file_with_notes)

pptx_file_with_notes
