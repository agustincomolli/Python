# Recrear la presentación para asegurar que no haya errores

from pptx import Presentation
from pptx.util import Inches

# Crear una nueva presentación
prs = Presentation()

# Definir las diapositivas y sus contenidos
slides_content = [
    ("El Sentido de lo Sagrado", "La Fe Católica frente a la Indiferencia Moderna"),
    ("Lo Tradicional y la Biblia", "La Iglesia vive de la Biblia y de la Tradición.\nImportancia igual de ambos según el Vaticano II (DV 9)."),
    ("La Iglesia y la Tradición", "La Biblia nace de la Tradición.\nSin la luz de la Tradición, la Biblia no nos valdría para nada."),
    ("Lo Sagrado", "Definición de lo Sagrado\nAlgo apartado o consagrado para un propósito espiritual.\nConsiderado digno de respeto o veneración especial."),
    ("Lo Sagrado vs. Lo Profano", "Distinción entre lo Sagrado y lo Profano\nLo sagrado está conectado con la divinidad.\nLo profano refiere a lo ordinario y cotidiano."),
    ("Manifestaciones de lo Sagrado", "Diversas Formas de lo Sagrado\n- Objetos Sagrados: libros religiosos, reliquias, etc.\n- Lugares Sagrados: templos, santuarios, etc.\n- Personas Sagradas: sacerdotes, monjes, etc.\n- Textos Sagrados: Biblia, Corán, etc.\n- Tiempo Sagrado: días santos, festivales, etc.\n- Acciones Sagradas: rituales, ceremonias, etc."),
    ("Continuidad de lo Sagrado", "De lo Pagano a lo Cristiano\nLa gracia perfecciona la naturaleza.\nContinuidad desde religiosidades naturales hasta la adoración cristiana."),
    ("Lo Sagrado Judío", "La Sacralidad en el Judaísmo\nOrden de sacralidades: fiestas, sacerdocio, lugares, etc.\nIsrael como pueblo sagrado."),
    ("Lo Sagrado Cristiano", "Cristo y la Sacralidad Cristiana\nCristo como sagrado absoluto.\nLa Iglesia como 'sacramento universal de salvación'."),
    ("Teología de lo Sagrado", "Definición Teológica\nJesucristo es sagrado por su humanidad.\nDistinción entre santo y sagrado."),
    ("Secularización y Desacralización", "Impacto Moderno en lo Sagrado\nPérdida de sensibilidad para lo sagrado en Occidente.\nConsecuencias espirituales de esta pérdida."),
    ("Espiritualidad Cristiana", "Amor a lo Sagrado\nImportancia del orden sacral en la espiritualidad católica.\nBúsqueda del Santo en las cosas sagradas de la Iglesia."),
    ("Tendencia de lo Sagrado a la Manifestación", "Visibilidad de lo Sagrado\nLo sagrado como signo claro y visible.\nEjemplos: templos, vestimentas religiosas, campanas, canto religioso."),
    ("Conclusión", "La Relevancia de lo Sagrado Hoy\nImportancia de mantener y respetar las formas sagradas.\nPapel de lo sagrado en la identidad y misión de la Iglesia.")
]

# Crear y añadir las diapositivas
for title, content in slides_content:
    slide_layout = prs.slide_layouts[1]  # Usar el layout de título y contenido
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

# Guardar la presentación en un archivo
pptx_file = "./Charla_sobre_Lo_Sagrado.pptx"
prs.save(pptx_file)

pptx_file
